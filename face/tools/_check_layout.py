from pptx import Presentation

p = Presentation(r'F:/Users/xiaoxuan/face/答辩/答辩PPT_人脸识别系统.pptx')
print(f'Slides: {len(p.slides)}')
all_text = []
for s in p.slides:
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    all_text.append(run.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            all_text.append(run.text)
full = ''.join(all_text)
print(f'Total chars: {len(full)}')
checks = ['Flask', 'dlib', 'ResNet', '128', 'L2', '欧氏距离', 'HOG', 'Algebraic', '2329',
         'xuan', '魏征', '谢宇', '王蕴X', '开发历程', '团队分工', '踩坑', 'Web 端']
for c in checks:
    mark = 'PASS' if c in full else 'MISS'
    print(f'  [{mark}] {c}')
n_pic = sum(1 for s in p.slides for sh in s.shapes if sh.shape_type == 13)
print(f'Embedded pictures: {n_pic}')
print('Image locations:')
for i, s in enumerate(p.slides, 1):
    for shape in s.shapes:
        if shape.shape_type == 13:
            x = shape.left/914400; y = shape.top/914400
            w = shape.width/914400; h = shape.height/914400
            print(f'  Slide {i:2d}: ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} ratio={w/h:.2f}')
