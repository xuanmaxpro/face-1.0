from pptx import Presentation

p = Presentation(r'F:/Users/xiaoxuan/face/答辩/答辩PPT_人脸识别系统.pptx')
print(f'Total slides: {len(p.slides)}')
print(f'Size: {p.slide_width/914400:.2f} x {p.slide_height/914400:.2f} inch')

# 抽取所有文本
all_text = []
for s in p.slides:
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    all_text.append(run.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            all_text.append(run.text)
full = ''.join(all_text)
print(f'Total chars: {len(full)}')
print()

checks = [
    'Flask', 'dlib', 'ResNet', '128', '0.6', 'L2', '欧氏距离', '68', '150', 'ResNet-34',
    'HOG', 'Pixel', 'Transform', 'Algebraic', '2329', 'EnsembleClassifier',
    'xuan', '魏征', '谢宇', '王蕴X', 'LFW',
    '开发历程', '团队分工', '踩坑', '测试验证', '总结', '致谢',
    'npz', 'json', 'source of truth', 'RectifiedAdam', 'nms',
]
for c in checks:
    status = 'OK ' if c in full else 'MISS'
    print(f'  [{status}] {c}')
print()
print('--- Page-by-page title check ---')
for i, s in enumerate(p.slides, 1):
    title = None
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt >= 20:
                        title = run.text
                        break
                if title:
                    break
        if title:
            break
    print(f'  Slide {i:2d}: {title}')
