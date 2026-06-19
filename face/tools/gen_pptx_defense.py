"""
生成答辩 PPT (20 张, 16:9, Pure Tech Blue 主题)
- 保留原 PPT 的核心内容 (背景/架构/算法/对比/总结)
- 新增开发历程/团队分工/踩坑与解决/演示 4 张 "开发经历" 主题页
- 输出: F:/Users/xiaoxuan/face/答辩/答辩PPT_人脸识别系统.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# ==================== 主题配色 (Pure Tech Blue) ====================
BG          = RGBColor(0xFF, 0xFF, 0xFF)   # 白底
PRIMARY     = RGBColor(0x03, 0x04, 0x5E)   # 深海军蓝 (标题/正文)
SECONDARY   = RGBColor(0x00, 0x77, 0xB6)   # 中蓝 (副标题/强调)
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)   # 亮蓝 (重点高亮)
LIGHT       = RGBColor(0xCA, 0xF0, 0xF8)   # 浅蓝 (卡片背景)
TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)   # 深灰 (正文)
TEXT_BODY   = RGBColor(0x37, 0x41, 0x51)   # 灰 (段落)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)   # 浅灰 (辅助)
DIVIDER     = RGBColor(0xE5, 0xE7, 0xEB)   # 分割线
SUCCESS     = RGBColor(0x10, 0xB9, 0x81)   # 绿 (通过/成功)
WARN        = RGBColor(0xF5, 0x9E, 0x0B)   # 橙 (警示)
DANGER      = RGBColor(0xEF, 0x44, 0x44)   # 红 (错误)

FONT_CN = "微软雅黑"
FONT_EN = "Arial"

OUT_PATH = "F:/Users/xiaoxuan/face/答辩/答辩PPT_人脸识别系统.pptx"
ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")


def add_image_with_frame(slide, x, y, w, h, image_path, frame_color=ACCENT, label=None):
    """
    插入图片 + 在底层叠加一个细边框矩形 (避免深色图片贴在白底上突兀)
    :param x, y, w, h: 图片位置和尺寸 (Inches)
    :param image_path: 图片绝对路径
    :param frame_color: 边框颜色, 默认 ACCENT
    :param label: 可选, 在图片下方加一行小标注
    """
    # 底层边框矩形 (略大于图片 0.04 边距)
    add_rect(slide, Inches(x - 0.04), Inches(y - 0.04), Inches(w + 0.08), Inches(h + 0.08),
             BG, line_color=frame_color, line_w=Pt(1.0))
    # 图片本体
    slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))
    # 标注
    if label:
        add_text(slide, Inches(x), Inches(y + h + 0.05), Inches(w), Inches(0.22),
                 label, font_size=9, font_color=TEXT_MUTED, align="center", italic=True)


# ==================== 底层工具 ====================
def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_w=None, shadow=False):
    """加一个矩形, 返回 shape"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_w is not None:
            shp.line.width = line_w
    if not shadow:
        # 关闭默认阴影
        sp = shp.shadow
        sp.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, font_size=14, font_color=TEXT_BODY,
             bold=False, align="left", anchor="top", font_cn=FONT_CN, font_en=FONT_EN,
             line_spacing=1.2, italic=False):
    """加文本框, 支持中英混排 (run 级别设中英字体)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # text 可能是 str, 也可能是 list of (text, options_dict)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }.get(align, PP_ALIGN.LEFT)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        # 处理多 run
        runs_spec = opts.get("runs", [(txt, {})])
        for j, (rt, ro) in enumerate(runs_spec):
            r = p.add_run() if j > 0 or p.runs else p.add_run()
            r.text = rt
            f = r.font
            f.name = font_en
            f.size = Pt(opts.get("font_size", font_size))
            f.bold = ro.get("bold", opts.get("bold", bold))
            f.italic = ro.get("italic", opts.get("italic", italic))
            # 中文字体 (rPr eastAsia)
            rPr = r._r.get_or_add_rPr()
            # 移除已存在 eastAsia
            for ea in rPr.findall(qn("a:ea")):
                rPr.remove(ea)
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", ro.get("font_cn", opts.get("font_cn", font_cn)))
            # 颜色
            color = ro.get("font_color", opts.get("font_color", font_color))
            f.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=ACCENT, width=1.5):
    """画一条线"""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


# ==================== 公共布局 ====================
def add_page_chrome(slide, page_num, total=20, section=""):
    """加页码 + 章节标识 (所有非封面页)"""
    # 顶部细线
    add_line(slide, Inches(0.5), Inches(0.42), Inches(9.5), Inches(0.42), DIVIDER, 0.5)
    # 底部页码
    add_text(slide, Inches(8.8), Inches(5.25), Inches(0.8), Inches(0.25),
             f"{page_num} / {total}", font_size=10, font_color=TEXT_MUTED,
             align="right")
    # 左下角: 项目标识
    if section:
        add_text(slide, Inches(0.5), Inches(5.25), Inches(6), Inches(0.25),
                 f"人脸识别系统  |  {section}", font_size=10, font_color=TEXT_MUTED)


def add_section_label(slide, x, y, label, color=ACCENT):
    """左上角小标签: '01 / 项目概述'"""
    # 短色块
    add_rect(slide, x, y, Inches(0.18), Inches(0.18), color)
    add_text(slide, x + Inches(0.25), y - Inches(0.04), Inches(3), Inches(0.3),
             label, font_size=11, font_color=color, bold=True, font_cn=FONT_CN)


def add_slide_title(slide, title, sub=None):
    """页面主标题 (大号 + 副标题)"""
    add_text(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.6),
             title, font_size=26, font_color=PRIMARY, bold=True)
    if sub:
        add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.35),
                 sub, font_size=13, font_color=TEXT_MUTED)
        # 标题下细装饰线
        add_line(slide, Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.5), ACCENT, 2.0)
    else:
        add_line(slide, Inches(0.5), Inches(1.15), Inches(1.5), Inches(1.15), ACCENT, 2.0)


# ==================== 各类 slide 模板 ====================
def make_cover(prs):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
    set_slide_bg(slide, BG)
    # 左侧蓝色色块
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), Inches(5.625), PRIMARY)
    # 顶部小标签
    add_text(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(0.3),
             "课程实践答辩  |  Program Design Practice",
             font_size=11, font_color=SECONDARY, bold=True)
    # 大标题
    add_text(slide, Inches(0.7), Inches(1.3), Inches(9), Inches(0.9),
             "基于 Flask + dlib 的人脸识别系统",
             font_size=36, font_color=PRIMARY, bold=True)
    # 副标题
    add_text(slide, Inches(0.7), Inches(2.25), Inches(9), Inches(0.6),
             "设计与实现  |  Design and Implementation",
             font_size=20, font_color=SECONDARY)
    # 装饰线
    add_line(slide, Inches(0.7), Inches(3.05), Inches(3.5), Inches(3.05), ACCENT, 3.0)
    # 关键标签
    add_text(slide, Inches(0.7), Inches(3.2), Inches(9), Inches(0.35),
             "1:1 确认  /  1:N 辨认  /  ResNet 128 维  /  4 种手工特征融合",
             font_size=14, font_color=TEXT_BODY)
    # 作者/老师
    add_text(slide, Inches(0.7), Inches(4.4), Inches(9), Inches(0.3),
             "答辩人: xuan, 魏征, 谢宇",
             font_size=13, font_color=TEXT_BODY, bold=True)
    add_text(slide, Inches(0.7), Inches(4.7), Inches(9), Inches(0.3),
             "指导教师: 王蕴X  副教授",
             font_size=12, font_color=TEXT_MUTED)
    add_text(slide, Inches(0.7), Inches(5.0), Inches(9), Inches(0.3),
             "2026 年 6 月",
             font_size=12, font_color=TEXT_MUTED)


def make_toc(prs):
    """Slide 2: 目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "CONTENTS  /  目录")
    add_slide_title(slide, "本次汇报的四个章节", "从项目背景到核心算法, 再到开发经历与未来展望")

    items = [
        ("01", "项目概述", "研究背景 / 核心功能 / 开发历程 / 团队分工", SECONDARY),
        ("02", "系统架构与核心算法", "5 层架构 / 8 步算法流程 / 预处理管线 / 特征匹配", ACCENT),
        ("03", "关键技术与对比实验", "4 种手工特征 / 深度 vs 手工 / 踩坑与解决 / 测试验证", PRIMARY),
        ("04", "总结与展望", "项目成果 / 不足反思 / 未来工作方向", SUCCESS),
    ]
    base_y = 1.85
    row_h = 0.78
    for i, (num, title, desc, color) in enumerate(items):
        y = Inches(base_y + i * row_h)
        # 数字块
        add_rect(slide, Inches(0.5), y, Inches(0.8), Inches(0.6), color)
        add_text(slide, Inches(0.5), y, Inches(0.8), Inches(0.6), num,
                 font_size=20, font_color=BG, bold=True, align="center", anchor="middle")
        # 标题
        add_text(slide, Inches(1.5), y + Inches(0.02), Inches(3.5), Inches(0.3),
                 title, font_size=16, font_color=PRIMARY, bold=True)
        # 描述
        add_text(slide, Inches(1.5), y + Inches(0.32), Inches(7), Inches(0.3),
                 desc, font_size=11, font_color=TEXT_MUTED)
    add_page_chrome(slide, 2, 20, "目录")


def make_section_divider(prs, num, title, intro, page_num):
    """章节封面: 大数字 + 标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    # 左侧大色块
    add_rect(slide, Inches(0), Inches(0), Inches(3.8), Inches(5.625), PRIMARY)
    # 大数字
    add_text(slide, Inches(0.5), Inches(1.4), Inches(3), Inches(2.0),
             num, font_size=140, font_color=ACCENT, bold=True, align="left")
    add_text(slide, Inches(0.5), Inches(3.3), Inches(3), Inches(0.4),
             "CHAPTER", font_size=12, font_color=LIGHT, bold=True)
    # 右侧标题
    add_text(slide, Inches(4.2), Inches(2.2), Inches(5.5), Inches(0.7),
             title, font_size=30, font_color=PRIMARY, bold=True)
    add_line(slide, Inches(4.2), Inches(3.0), Inches(5.5), Inches(3.0), ACCENT, 3.0)
    add_text(slide, Inches(4.2), Inches(3.15), Inches(5.5), Inches(1.2),
             intro, font_size=14, font_color=TEXT_BODY, line_spacing=1.5)
    add_page_chrome(slide, page_num, 20, f"第{num}章")


# ==================== 第 1 章: 项目概述 ====================
def make_bg_meaning(prs, page_num):
    """Slide 4: 研究背景与意义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "01  /  项目概述")
    add_slide_title(slide, "研究背景与意义", "人脸识别是计算机视觉的经典任务, 也是本门课程实践的核心选题")

    # 左侧: 现状 (卡片)
    add_rect(slide, Inches(0.5), Inches(1.75), Inches(4.5), Inches(3.3), LIGHT)
    add_text(slide, Inches(0.7), Inches(1.85), Inches(4.2), Inches(0.4),
             "人脸识别技术的现状", font_size=15, font_color=PRIMARY, bold=True)
    add_text(slide, Inches(0.7), Inches(2.3), Inches(4.2), Inches(2.6),
             [
                 ("核心目标: 让计算机自动识别与验证人脸身份", {"font_size": 12, "font_color": TEXT_BODY, "line_spacing": 1.4}),
                 ("", {}),
                 ("典型应用场景:", {"font_size": 12, "font_color": PRIMARY, "bold": True, "line_spacing": 1.4}),
                 ("  • 安防监控 / 门禁系统", {"font_size": 11, "font_color": TEXT_BODY, "line_spacing": 1.4}),
                 ("  • 身份认证 / 手机解锁 / 课堂签到", {"font_size": 11, "font_color": TEXT_BODY, "line_spacing": 1.4}),
                 ("", {}),
                 ("两大主流技术路线:", {"font_size": 12, "font_color": PRIMARY, "bold": True, "line_spacing": 1.4}),
                 ("  • 传统手工特征: HOG, LBP, SIFT (可解释强)", {"font_size": 11, "font_color": TEXT_BODY, "line_spacing": 1.4}),
                 ("  • 深度学习: ResNet, FaceNet (精度高, 黑盒强)", {"font_size": 11, "font_color": TEXT_BODY, "line_spacing": 1.4}),
             ])

    # 右侧: 本项目研究价值
    add_rect(slide, Inches(5.2), Inches(1.75), Inches(4.3), Inches(3.3), BG,
             line_color=ACCENT, line_w=Pt(1.5))
    add_text(slide, Inches(5.4), Inches(1.85), Inches(4), Inches(0.4),
             "本项目的研究价值", font_size=15, font_color=PRIMARY, bold=True)
    values = [
        ("① 端到端 Web 化", "Flask + dlib 构建可用的人脸识别 Web 应用, 支持在线录入与实时识别"),
        ("② 双模式识别", "同时实现 1:1 确认 (验证) 和 1:N 辨认 (识别) 两种典型场景"),
        ("③ 特征对比实验", "把 dlib 深度学习特征与 4 种手工特征融合进行横向对比, 验证技术路线差异"),
        ("④ 工程化思考", "统一预处理管线 + 数据一致性同步, 体现工程化设计意识"),
    ]
    for i, (head, body) in enumerate(values):
        y = Inches(2.3 + i * 0.65)
        add_text(slide, Inches(5.4), y, Inches(4), Inches(0.3),
                 head, font_size=12, font_color=ACCENT, bold=True)
        add_text(slide, Inches(5.4), y + Inches(0.28), Inches(4), Inches(0.35),
                 body, font_size=10, font_color=TEXT_BODY)

    # 底部技术栈
    add_rect(slide, Inches(0.5), Inches(5.05), Inches(9), Inches(0.16), LIGHT)
    add_text(slide, Inches(0.6), Inches(5.07), Inches(9), Inches(0.14),
             "技术栈:  Python 3.11  |  Flask 3.0  |  dlib 19.24  |  OpenCV 4.8  |  NumPy 1.24  |  PptxGenJS / python-pptx",
             font_size=9, font_color=PRIMARY, anchor="middle")
    add_page_chrome(slide, page_num, 20, "研究背景与意义")


def make_features(prs, page_num):
    """Slide 5: 核心功能与亮点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "01  /  项目概述")
    add_slide_title(slide, "核心功能与亮点", "系统支持的两大核心识别模式 + 13 个 API 端点")

    # 顶部: 两种识别模式对比
    add_text(slide, Inches(0.5), Inches(1.7), Inches(4.5), Inches(0.3),
             "两种核心识别模式", font_size=14, font_color=PRIMARY, bold=True)

    # 左卡: 1:1 确认
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(2.1), Inches(1.4), PRIMARY)
    add_text(slide, Inches(0.5), Inches(2.1), Inches(2.1), Inches(0.4),
             "1 : 1  确认", font_size=18, font_color=BG, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(2.55), Inches(2.1), Inches(0.35),
             "Verification", font_size=10, font_color=LIGHT, align="center")
    add_text(slide, Inches(0.5), Inches(2.95), Inches(2.1), Inches(0.5),
             "是 这个人 吗?",
             font_size=12, font_color=BG, align="center", anchor="middle")

    # 右卡: 1:N 辨认
    add_rect(slide, Inches(2.7), Inches(2.05), Inches(2.1), Inches(1.4), SECONDARY)
    add_text(slide, Inches(2.7), Inches(2.1), Inches(2.1), Inches(0.4),
             "1 : N  辨认", font_size=18, font_color=BG, bold=True, align="center")
    add_text(slide, Inches(2.7), Inches(2.55), Inches(2.1), Inches(0.35),
             "Identification", font_size=10, font_color=LIGHT, align="center")
    add_text(slide, Inches(2.7), Inches(2.95), Inches(2.1), Inches(0.5),
             "这是 谁 ?",
             font_size=12, font_color=BG, align="center", anchor="middle")

    # 右侧: 系统能力描述
    add_rect(slide, Inches(5.0), Inches(2.05), Inches(4.5), Inches(1.4), LIGHT)
    add_text(slide, Inches(5.2), Inches(2.15), Inches(4.2), Inches(0.3),
             "系统能力速览", font_size=12, font_color=PRIMARY, bold=True)
    bullets = [
        "Web 上传图片 / 摄像头实时识别",
        "支持多人脸同时检测与识别",
        "13 个 RESTful API 端点, 模块解耦",
        "LFW 公开数据集 + 个人照片演示",
    ]
    for i, b in enumerate(bullets):
        add_text(slide, Inches(5.2), Inches(2.5 + i * 0.22), Inches(4.2), Inches(0.22),
                 f"  •  {b}", font_size=10, font_color=TEXT_BODY)

    # 底部: 系统 Web 端实测截图 (两种识别模式, 保持 2:1 比例, 避开顶部卡 2.05-3.45)
    # 2 张图并排: W=2.8 x H=1.4 保持 2:1, 居中, 总宽 6.0
    img_w, img_h, gap = 2.8, 1.4, 0.4
    total_w = 2 * img_w + gap  # 6.0
    left_x = (10 - total_w) / 2  # 2.0
    img_y = 3.55
    # 左图: 上传识别
    add_image_with_frame(
        slide, x=left_x, y=img_y, w=img_w, h=img_h,
        image_path=os.path.join(ASSETS_DIR, "upload_recognize.png"),
        frame_color=PRIMARY,
        label="上传识别  -  1:1 / 1:N 双模式"
    )
    # 右图: 实时识别
    add_image_with_frame(
        slide, x=left_x + img_w + gap, y=img_y, w=img_w, h=img_h,
        image_path=os.path.join(ASSETS_DIR, "realtime_recognize.png"),
        frame_color=ACCENT,
        label="实时识别  -  摄像头 + 设置 + 记录"
    )

    add_page_chrome(slide, page_num, 20, "核心功能与亮点")


def make_timeline(prs, page_num):
    """Slide 6: 开发历程 (时间线) - 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "01  /  项目概述")
    add_slide_title(slide, "开发历程", "从 0 到 1: 一段不断踩坑和迭代的过程")

    # 横向时间线
    line_y = 2.7
    add_line(slide, Inches(0.8), Inches(line_y), Inches(9.2), Inches(line_y), ACCENT, 2.0)

    phases = [
        ("第 1 周", "需求 + 调研", "读任务书, 选型 dlib / Flask, 搭项目骨架"),
        ("第 2 周", "核心算法", "打通 检测 → 关键点 → ResNet 128 维, 本地 CLI 跑通"),
        ("第 3 周", "Web 化", "Flask 入口 + 13 个 API + index.html 前端"),
        ("第 4 周", "4 种手工特征", "HOG / Pixel / Transform / Algebraic + Ensemble"),
        ("第 5 周", "踩坑 + 重构", "L2 归一化 bug fix, npz/json 同步重构, 7 个测试"),
        ("第 6 周", "完善 + 答辩", "文档 / 讲稿 / PPT, 准备现场演示"),
    ]
    n = len(phases)
    span = 8.4  # 0.8 ~ 9.2
    step = span / (n - 1)
    for i, (week, head, body) in enumerate(phases):
        cx = 0.8 + i * step
        # 圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.13), Inches(line_y - 0.13), Inches(0.26), Inches(0.26))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.color.rgb = BG; dot.line.width = Pt(2)
        # 周次 (在圆点上方)
        add_text(slide, Inches(cx - 0.6), Inches(line_y - 0.55), Inches(1.2), Inches(0.3),
                 week, font_size=10, font_color=ACCENT, bold=True, align="center")
        # 阶段名 (在圆点下方)
        add_text(slide, Inches(cx - 0.7), Inches(line_y + 0.2), Inches(1.4), Inches(0.3),
                 head, font_size=12, font_color=PRIMARY, bold=True, align="center")
        # 描述 (再下)
        add_text(slide, Inches(cx - 0.85), Inches(line_y + 0.55), Inches(1.7), Inches(0.9),
                 body, font_size=9, font_color=TEXT_BODY, align="center", line_spacing=1.3)

    # 底部: 心得一句话
    add_rect(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.55), LIGHT)
    add_text(slide, Inches(0.7), Inches(4.65), Inches(8.6), Inches(0.45),
             "心得: \"做出来\" 比 \"做完美\" 更重要 — 先跑通, 再迭代; 每一行 bug 都是进步.",
             font_size=12, font_color=PRIMARY, bold=True, anchor="middle")
    add_page_chrome(slide, page_num, 20, "开发历程")


def make_team(prs, page_num):
    """Slide 7: 团队分工 (3 列卡片) - 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "01  /  项目概述")
    add_slide_title(slide, "团队分工", "3 人小组成员各司其职, 每周碰头对齐进度")

    members = [
        ("xuan", "项目负责人 / 核心算法",
         [
             "项目整体架构设计",
             "dlib 路径核心代码 (face_recognition_core.py)",
             "数据一致性 (npz + json 同步) 重构",
             "答辩 PPT 与讲稿主笔",
         ], PRIMARY),
        ("魏征", "Web 前端 / Flask 后端",
         [
             "Flask 入口与 13 个 API 端点",
             "前端 index.html / style.css / main.js",
             "摄像头实时帧处理 (节流 / 跳帧)",
             "UI 交互与提示信息完善",
         ], SECONDARY),
        ("谢宇", "4 种手工特征 / 对比实验",
         [
             "HOG / Pixel / Transform / Algebraic 实现",
             "EnsembleClassifier 加权融合",
             "ensemble_database.npz 独立库",
             "对比实验设计与文档撰写",
         ], ACCENT),
    ]
    for i, (name, role, tasks, color) in enumerate(members):
        x = Inches(0.5 + i * 3.05)
        # 卡片背景
        add_rect(slide, x, Inches(1.75), Inches(2.9), Inches(3.2), BG,
                 line_color=color, line_w=Pt(1.5))
        # 顶部色块
        add_rect(slide, x, Inches(1.75), Inches(2.9), Inches(0.7), color)
        # 姓名
        add_text(slide, x, Inches(1.78), Inches(2.9), Inches(0.4),
                 name, font_size=18, font_color=BG, bold=True, align="center", anchor="middle")
        # 角色
        add_text(slide, x, Inches(2.18), Inches(2.9), Inches(0.25),
                 role, font_size=10, font_color=LIGHT, align="center")
        # 任务列表
        for j, task in enumerate(tasks):
            ty = Inches(2.6 + j * 0.5)
            # 圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), ty + Inches(0.1), Inches(0.08), Inches(0.08))
            dot.fill.solid(); dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            add_text(slide, x + Inches(0.3), ty, Inches(2.5), Inches(0.45),
                     task, font_size=10, font_color=TEXT_BODY, line_spacing=1.3)
    add_page_chrome(slide, page_num, 20, "团队分工")


# ==================== 第 2 章: 系统架构与核心算法 ====================
def make_architecture(prs, page_num):
    """Slide 9: 系统整体架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "02  /  系统架构与核心算法")
    add_slide_title(slide, "系统整体架构", "5 层模块化设计, 每层职责单一, 接口清晰")

    # 5 层架构 (从下到上, 视觉上从上到下展示请求流)
    layers = [
        ("前端层 (Browser)", "index.html + main.js", "图片上传 / 摄像头 / 结果绘制 / 人员管理", ACCENT),
        ("Web 层 (Flask App)", "app.py  (449 行)", "13 个 RESTful API / 请求路由 / 静态资源", SECONDARY),
        ("核心层 (FaceRecognitionCore)", "face_recognition_core.py  (510 行)", "人脸检测 / 关键点 / 特征提取 / 1:1 与 1:N 匹配", PRIMARY),
        ("特征层 (FeatureExtractor)", "feature_extractor.py  (286 行)", "HOG / Pixel / Transform / Algebraic + Ensemble", SUCCESS),
        ("数据层 (Database)", "database.py  (134 行) + data/", "FaceDatabase / RecognitionRecord / npz + json 落盘", WARN),
    ]
    base_y = 1.8
    row_h = 0.62
    for i, (head, file_, desc, color) in enumerate(layers):
        y = Inches(base_y + i * row_h)
        add_rect(slide, Inches(0.5), y, Inches(9), Inches(0.5), color)
        # 标题
        add_text(slide, Inches(0.7), y, Inches(3.0), Inches(0.5),
                 head, font_size=12, font_color=BG, bold=True, anchor="middle")
        # 文件
        add_text(slide, Inches(3.8), y, Inches(2.5), Inches(0.5),
                 file_, font_size=10, font_color=BG, anchor="middle", italic=True)
        # 描述
        add_text(slide, Inches(6.4), y, Inches(3.0), Inches(0.5),
                 desc, font_size=10, font_color=BG, anchor="middle")
        # 箭头 (除了最后一层)
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.85), y + Inches(0.51), Inches(0.3), Inches(0.09))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()

    # 底部: 一句话总结
    add_text(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
             "设计原则: 单一职责 + 数据一致 (npz 单一 source of truth) + 接口解耦 (13 个 API 端点)",
             font_size=11, font_color=PRIMARY, bold=True, align="center")
    add_page_chrome(slide, page_num, 20, "系统整体架构")


def make_algorithm_flow(prs, page_num):
    """Slide 10: 人脸识别算法流程 (8 步流程图)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "02  /  系统架构与核心算法")
    add_slide_title(slide, "人脸识别算法流程", "8 个步骤把一张人脸图变成一个 128 维数字, 再跟库比对")

    steps = [
        ("①", "输入图片", "用户上传 / 摄像头帧"),
        ("②", "灰度化", "BGR → Gray"),
        ("③", "人脸检测", "HOG + SVM 滑窗"),
        ("④", "68 点关键点", "ERT 级联回归"),
        ("⑤", "几何对齐", "dlib.get_face_chip\n→ 150×150"),
        ("⑥", "预处理", "直方图 + 高斯"),
        ("⑦", "ResNet 推理", "→ 128 维向量"),
        ("⑧", "L2 归一化 + 比对", "欧氏距离 < 0.6 判同一人"),
    ]
    n = len(steps)
    span = 8.6
    base_x = 0.7
    step_w = 1.05
    gap = (span - n * step_w) / (n - 1) if n > 1 else 0
    flow_y = 2.2
    box_h = 0.85
    for i, (num, head, body) in enumerate(steps):
        x = Inches(base_x + i * (step_w + gap))
        # 数字标
        add_rect(slide, x, Inches(flow_y), Inches(step_w), Inches(0.25), ACCENT)
        add_text(slide, x, Inches(flow_y), Inches(step_w), Inches(0.25),
                 num, font_size=10, font_color=BG, bold=True, align="center", anchor="middle")
        # 标题
        add_rect(slide, x, Inches(flow_y + 0.25), Inches(step_w), Inches(0.32), PRIMARY)
        add_text(slide, x, Inches(flow_y + 0.25), Inches(step_w), Inches(0.32),
                 head, font_size=10, font_color=BG, bold=True, align="center", anchor="middle")
        # 描述
        add_rect(slide, x, Inches(flow_y + 0.57), Inches(step_w), Inches(0.42), LIGHT)
        add_text(slide, x, Inches(flow_y + 0.6), Inches(step_w), Inches(0.42),
                 body, font_size=8, font_color=PRIMARY, align="center", anchor="middle", line_spacing=1.2)
        # 箭头 (除了最后)
        if i < n - 1:
            ax = base_x + i * (step_w + gap) + step_w
            ay = flow_y + 0.32
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax + 0.02), Inches(ay - 0.05), Inches(gap - 0.04), Inches(0.1))
            arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()

    # 关键点说明 (3 段并排)
    add_text(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.3),
             "关键点解析", font_size=14, font_color=PRIMARY, bold=True)
    notes = [
        ("dlib HOG+SVM 检测", "比 CNN 轻量, CPU 实时, 对正面人脸检出率高"),
        ("ERT 68 点关键点", "级联回归树, 比传统 ASM 更快更准"),
        ("ResNet 128 维", "在 300 万张人脸上预训练, 已是工业级精度"),
    ]
    for i, (head, body) in enumerate(notes):
        x = Inches(0.5 + i * 3.05)
        add_rect(slide, x, Inches(3.95), Inches(2.9), Inches(1.05), BG,
                 line_color=ACCENT, line_w=Pt(1.0))
        add_text(slide, x + Inches(0.15), Inches(4.05), Inches(2.6), Inches(0.3),
                 head, font_size=11, font_color=ACCENT, bold=True)
        add_text(slide, x + Inches(0.15), Inches(4.4), Inches(2.6), Inches(0.55),
                 body, font_size=10, font_color=TEXT_BODY, line_spacing=1.3)
    add_page_chrome(slide, page_num, 20, "人脸识别算法流程")


def make_preprocess(prs, page_num):
    """Slide 11: 图像预处理管线 (7 原子 + 2 管线)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "02  /  系统架构与核心算法")
    add_slide_title(slide, "图像预处理管线", "7 个独立原子方法, 通过 apply_preprocess_pipeline 统一调度")

    # 上半部分: 7 个原子方法 (4 + 3 布局)
    add_text(slide, Inches(0.5), Inches(1.7), Inches(9), Inches(0.3),
             "7 个原子预处理方法", font_size=14, font_color=PRIMARY, bold=True)
    atoms = [
        ("光线补偿", "LAB 空间亮度调整"),
        ("灰度变换", "线性亮度拉伸"),
        ("直方图均衡", "YUV 通道均衡"),
        ("归一化", "像素值归一到 [0,255]"),
        ("几何校正", "dlib.get_face_chip"),
        ("高斯滤波", "去除图像噪声"),
        ("锐化", "卷积核增强边缘"),
    ]
    # 4 个一行
    for i, (head, body) in enumerate(atoms[:4]):
        x = Inches(0.5 + i * 2.27)
        add_rect(slide, x, Inches(2.05), Inches(2.15), Inches(0.6), LIGHT)
        add_text(slide, x, Inches(2.08), Inches(2.15), Inches(0.25),
                 head, font_size=10, font_color=PRIMARY, bold=True, align="center")
        add_text(slide, x, Inches(2.32), Inches(2.15), Inches(0.3),
                 body, font_size=8, font_color=TEXT_BODY, align="center")
    # 3 个一行 (居中)
    offset = (9 - 3 * 2.27) / 2
    for i, (head, body) in enumerate(atoms[4:]):
        x = Inches(offset + i * 2.27)
        add_rect(slide, x, Inches(2.7), Inches(2.15), Inches(0.6), LIGHT)
        add_text(slide, x, Inches(2.73), Inches(2.15), Inches(0.25),
                 head, font_size=10, font_color=PRIMARY, bold=True, align="center")
        add_text(slide, x, Inches(2.97), Inches(2.15), Inches(0.3),
                 body, font_size=8, font_color=TEXT_BODY, align="center")

    # 下半部分: 两条管线对比
    add_text(slide, Inches(0.5), Inches(3.55), Inches(9), Inches(0.3),
             "两条预配置管线 (按场景选用)", font_size=14, font_color=PRIMARY, bold=True)
    # 完整管线
    add_rect(slide, Inches(0.5), Inches(3.9), Inches(4.4), Inches(1.25), BG,
             line_color=PRIMARY, line_w=Pt(1.0))
    add_rect(slide, Inches(0.5), Inches(3.9), Inches(4.4), Inches(0.3), PRIMARY)
    add_text(slide, Inches(0.5), Inches(3.9), Inches(4.4), Inches(0.3),
             "PREPROCESS_PIPELINE_FULL  -  录入场景 (7 步)", font_size=10, font_color=BG, bold=True, align="center", anchor="middle")
    add_text(slide, Inches(0.65), Inches(4.3), Inches(4.1), Inches(0.8),
             "光线补偿 → 几何校正 → 直方图均衡 → 归一化 → 滤波 → 锐化 → 直方图均衡",
             font_size=10, font_color=TEXT_BODY, line_spacing=1.5, anchor="middle")
    # 快速管线
    add_rect(slide, Inches(5.1), Inches(3.9), Inches(4.4), Inches(1.25), BG,
             line_color=ACCENT, line_w=Pt(1.0))
    add_rect(slide, Inches(5.1), Inches(3.9), Inches(4.4), Inches(0.3), ACCENT)
    add_text(slide, Inches(5.1), Inches(3.9), Inches(4.4), Inches(0.3),
             "PREPROCESS_PIPELINE_FAST  -  识别场景 (3 步)", font_size=10, font_color=BG, bold=True, align="center", anchor="middle")
    add_text(slide, Inches(5.25), Inches(4.3), Inches(4.1), Inches(0.8),
             "对齐 (align) → 直方图均衡 → 高斯滤波",
             font_size=10, font_color=TEXT_BODY, line_spacing=1.5, anchor="middle")
    add_page_chrome(slide, page_num, 20, "图像预处理管线")


def make_feature_match(prs, page_num):
    """Slide 12: 特征提取与匹配"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "02  /  系统架构与核心算法")
    add_slide_title(slide, "特征提取与匹配", "128 维向量 + 欧氏距离 + 阈值判定")

    # 左: 特征提取
    add_text(slide, Inches(0.5), Inches(1.7), Inches(4.5), Inches(0.3),
             "特征提取 (深度学习路径)", font_size=14, font_color=PRIMARY, bold=True)
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(4.5), Inches(2.6), LIGHT)
    feat_lines = [
        ("• 模型", "dlib ResNet-34 预训练 (300 万人脸)"),
        ("• 输入", "150×150 对齐后的 RGB 人脸图"),
        ("• 输出", "128 维浮点特征向量"),
        ("• 归一化", "L2 归一化, 让向量落在单位球上"),
        ("• 核心代码", "face_model.compute_face_descriptor()"),
    ]
    for i, (k, v) in enumerate(feat_lines):
        y = Inches(2.18 + i * 0.45)
        add_text(slide, Inches(0.7), y, Inches(1.0), Inches(0.35),
                 k, font_size=11, font_color=ACCENT, bold=True)
        add_text(slide, Inches(1.7), y, Inches(3.2), Inches(0.4),
                 v, font_size=11, font_color=TEXT_BODY)

    # 右: 匹配算法
    add_text(slide, Inches(5.2), Inches(1.7), Inches(4.5), Inches(0.3),
             "匹配与识别", font_size=14, font_color=PRIMARY, bold=True)
    add_rect(slide, Inches(5.2), Inches(2.05), Inches(4.3), Inches(2.6), BG,
             line_color=ACCENT, line_w=Pt(1.0))
    # 公式 1: 距离
    add_text(slide, Inches(5.4), Inches(2.2), Inches(4), Inches(0.3),
             "距离度量 (欧氏距离)", font_size=11, font_color=PRIMARY, bold=True)
    add_text(slide, Inches(5.4), Inches(2.5), Inches(4), Inches(0.4),
             "d = || a - b ||₂  =  √(Σ (aᵢ - bᵢ)²)",
             font_size=14, font_color=ACCENT, bold=True, font_en="Cambria")
    # 公式 2: 相似度
    add_text(slide, Inches(5.4), Inches(2.95), Inches(4), Inches(0.3),
             "相似度 (距离 → 相似度, clamp 到 [0,1])", font_size=11, font_color=PRIMARY, bold=True)
    add_text(slide, Inches(5.4), Inches(3.25), Inches(4), Inches(0.4),
             "similarity = max(0, 1 - d)",
             font_size=14, font_color=ACCENT, bold=True, font_en="Cambria")
    # 阈值
    add_text(slide, Inches(5.4), Inches(3.7), Inches(4), Inches(0.3),
             "判定阈值", font_size=11, font_color=PRIMARY, bold=True)
    add_text(slide, Inches(5.4), Inches(4.0), Inches(4), Inches(0.4),
             "d < 0.6  →  判为同一人  (dlib 官方推荐)",
             font_size=13, font_color=SUCCESS, bold=True)
    # 两种模式
    add_text(slide, Inches(5.4), Inches(4.35), Inches(4), Inches(0.25),
             "两种工作模式", font_size=11, font_color=PRIMARY, bold=True)
    add_text(slide, Inches(5.4), Inches(4.6), Inches(4), Inches(0.3),
             "1:1 确认  verify_face()  /  1:N 辨认  identify_face()",
             font_size=10, font_color=TEXT_BODY)

    # 底部: 录入时取平均
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4), PRIMARY)
    add_text(slide, Inches(0.7), Inches(4.83), Inches(8.6), Inches(0.35),
             "录入增强: 对一个人的多张图都提特征, 取平均再 L2 归一化, 比单张图更鲁棒",
             font_size=12, font_color=BG, bold=True, anchor="middle")
    add_page_chrome(slide, page_num, 20, "特征提取与匹配")


# ==================== 第 3 章: 关键技术与对比实验 ====================
def make_handcraft(prs, page_num):
    """Slide 14: 4 种手工特征 (表格)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "03  /  关键技术与对比实验")
    add_slide_title(slide, "4 种手工特征提取方法", "与传统机器学习思路对齐, 跟深度学习做横向对比")

    # 表格
    headers = ["特征类型", "维度", "实现方法", "形象解释"]
    rows = [
        ["HOG", "1764", "cv2.HOGDescriptor", "梯度方向直方图, 看人脸的纹路走向"],
        ["Pixel Stats", "14", "统计量 + 直方图", "看人脸的明暗分布 / 颜色直方图"],
        ["Transform", "528", "DCT + PCA + DFT", "看人脸的频率特征 / 主成分"],
        ["Algebraic", "23", "SVD + 范数 + LBP", "看人脸的矩阵代数属性"],
        ["Ensemble", "2329", "4 种 concat + 加权余弦", "4 种融合, 跟 dlib 对比实验"],
    ]
    table_x = Inches(0.5); table_y = Inches(1.7); table_w = Inches(9); table_h = Inches(1.7)
    table = slide.shapes.add_table(len(rows) + 1, 4, table_x, table_y, table_w, table_h).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(1.0)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(4.3)
    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = BG
        r.font.name = FONT_EN
        rPr = r._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn("a:ea")); ea.set("typeface", FONT_CN)
    # 数据行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            if ci == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(10)
            r.font.color.rgb = PRIMARY if ci == 0 else TEXT_BODY
            r.font.bold = (ci == 0)
            r.font.name = FONT_EN
            rPr = r._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn("a:ea")); ea.set("typeface", FONT_CN)
            # 行背景
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 == 0 else BG

    # 加权说明 + 系统截图 (两栏布局: 左侧加权小卡横排, 右侧截图)
    add_text(slide, Inches(0.5), Inches(3.55), Inches(9), Inches(0.3),
             "EnsembleClassifier 加权策略  +  系统 Web 端实测", font_size=13, font_color=PRIMARY, bold=True)
    # 左侧: 加权小卡 (1 行 4 个, 高度 0.5", X 范围 0.5-6.5)
    weights = [
        ("HOG", "0.40", "梯度方向最重要", ACCENT),
        ("Transform", "0.25", "频域特征次之", SECONDARY),
        ("Algebraic", "0.25", "代数属性补充", PRIMARY),
        ("Pixel Stats", "0.10", "像素信息量较少", TEXT_MUTED),
    ]
    for i, (name, w, desc, color) in enumerate(weights):
        x = Inches(0.5 + i * 1.55)
        add_rect(slide, x, Inches(3.95), Inches(1.45), Inches(0.5), BG, line_color=color, line_w=Pt(0.75))
        add_text(slide, x + Inches(0.08), Inches(3.98), Inches(0.7), Inches(0.22),
                 name, font_size=9, font_color=color, bold=True)
        add_text(slide, x + Inches(0.8), Inches(3.98), Inches(0.55), Inches(0.22),
                 f"x {w}", font_size=10, font_color=color, bold=True, align="right", font_en="Cambria")
        add_text(slide, x + Inches(0.08), Inches(4.18), Inches(1.3), Inches(0.3),
                 desc, font_size=7, font_color=TEXT_MUTED, line_spacing=1.2)

    # 右侧: 系统 Web 端特征介绍截图 (2:1 比例, W=2.0 x H=1.0)
    add_image_with_frame(
        slide, x=6.8, y=3.95, w=2.0, h=1.0,
        image_path=os.path.join(ASSETS_DIR, "feature_info.png"),
        frame_color=ACCENT,
        label="系统特征选择界面 (Web 端)"
    )
    add_page_chrome(slide, page_num, 20, "4 种手工特征")


def make_comparison(prs, page_num):
    """Slide 15: 深度学习 vs 手工特征 (对比)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "03  /  关键技术与对比实验")
    add_slide_title(slide, "深度学习 vs 手工特征", "两条技术路线的横向对比")

    # 左右两栏
    add_text(slide, Inches(0.5), Inches(1.7), Inches(4.5), Inches(0.3),
             "ResNet 深度学习 (主路径)", font_size=14, font_color=PRIMARY, bold=True, align="center")
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(4.5), Inches(2.6), BG,
             line_color=PRIMARY, line_w=Pt(1.5))
    dl_items = [
        ("特征维度", "128 维"),
        ("模型基础", "ResNet-34 残差网络"),
        ("训练数据", "dlib 官方 300 万人脸预训练"),
        ("归一化", "L2 归一化到单位向量"),
        ("距离度量", "欧氏距离 d = ||a-b||₂"),
        ("阈值", "d < 0.6 判同一人"),
        ("权重", "端到端学习, 无人工权重"),
    ]
    for i, (k, v) in enumerate(dl_items):
        y = Inches(2.18 + i * 0.33)
        add_text(slide, Inches(0.7), y, Inches(1.4), Inches(0.3),
                 k, font_size=10, font_color=PRIMARY, bold=True)
        add_text(slide, Inches(2.1), y, Inches(2.8), Inches(0.3),
                 v, font_size=10, font_color=TEXT_BODY)

    add_text(slide, Inches(5.2), Inches(1.7), Inches(4.5), Inches(0.3),
             "Ensemble 手工特征 (对比实验)", font_size=14, font_color=ACCENT, bold=True, align="center")
    add_rect(slide, Inches(5.2), Inches(2.05), Inches(4.3), Inches(2.6), BG,
             line_color=ACCENT, line_w=Pt(1.5))
    hc_items = [
        ("特征维度", "2329 维 (HOG+Pixel+Trans+Alg)"),
        ("模型基础", "经典图像处理算子组合"),
        ("训练数据", "无需训练, 直接从图片提取"),
        ("归一化", "各特征独立归一化"),
        ("距离度量", "加权余弦相似度"),
        ("阈值", "余弦相似度 > 0.6 判同一人"),
        ("权重", "HOG 0.4 / Trans 0.25 / Alg 0.25 / Pixel 0.1"),
    ]
    for i, (k, v) in enumerate(hc_items):
        y = Inches(2.18 + i * 0.33)
        add_text(slide, Inches(5.4), y, Inches(1.4), Inches(0.3),
                 k, font_size=10, font_color=ACCENT, bold=True)
        add_text(slide, Inches(6.8), y, Inches(2.6), Inches(0.3),
                 v, font_size=10, font_color=TEXT_BODY)

    # 底部: 演示方式
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4), PRIMARY)
    add_text(slide, Inches(0.7), Inches(4.83), Inches(8.6), Inches(0.35),
             "现场演示: 同一张测试图, 调 /api/recognize (ResNet) 和 /api/ensemble/recognize (手工) 对比识别结果",
             font_size=12, font_color=BG, bold=True, anchor="middle")
    add_page_chrome(slide, page_num, 20, "深度学习 vs 手工特征")


def make_pitfalls(prs, page_num):
    """Slide 16: 踩坑与解决 (3 张卡片) - 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "03  /  关键技术与对比实验")
    add_slide_title(slide, "踩坑与解决", "开发过程中遇到的 3 个真实问题, 以及如何修复")

    pitfalls = [
        ("Bug 1: 欧氏距离虚高", "录入和识别未做 L2 归一化, 光线亮的人脸距离数值偏大, 容易误判",
         "在 extract_features 末尾增加 _l2_normalize():  vec / ||vec||, 让所有向量长度都等于 1, 距离范围稳定在 [0, 2] 之间可比",
         DANGER, "face_recognition_core.py:188  _l2_normalize"),
        ("Bug 2: 相似度出现负值", "原公式 1 - distance 在归一化前会算出负数 (反向量场景), 前端显示 -50% 误导用户",
         "在 _distance_to_similarity 中加 max(0, 1-d) clamp 到 [0, 1], 负值压到 0, 物理上更合理",
         WARN, "face_recognition_core.py:194  _distance_to_similarity"),
        ("重构: 双文件数据不一致", "face_db.json 和 face_database.npz 各存一份 features, 录入时只更新一边, 容易错位",
         "改为 npz 单一 source of truth, json 只存元数据; 每次 save_database 自动 sync_metadata_from_npz, 启动时自动重建",
         SUCCESS, "face_recognition_core.py:391  save_database"),
    ]
    for i, (title, sym, fix, color, code_ref) in enumerate(pitfalls):
        x = Inches(0.5 + i * 3.05)
        # 卡片
        add_rect(slide, x, Inches(1.75), Inches(2.9), Inches(3.2), BG, line_color=color, line_w=Pt(1.5))
        # 顶部色块
        add_rect(slide, x, Inches(1.75), Inches(2.9), Inches(0.45), color)
        add_text(slide, x, Inches(1.75), Inches(2.9), Inches(0.45),
                 title, font_size=12, font_color=BG, bold=True, align="center", anchor="middle")
        # 现象
        add_text(slide, x + Inches(0.15), Inches(2.3), Inches(2.6), Inches(0.25),
                 "现象", font_size=10, font_color=color, bold=True)
        add_text(slide, x + Inches(0.15), Inches(2.55), Inches(2.6), Inches(0.85),
                 sym, font_size=9, font_color=TEXT_BODY, line_spacing=1.3)
        # 解决
        add_text(slide, x + Inches(0.15), Inches(3.4), Inches(2.6), Inches(0.25),
                 "解决方案", font_size=10, font_color=color, bold=True)
        add_text(slide, x + Inches(0.15), Inches(3.65), Inches(2.6), Inches(1.0),
                 fix, font_size=9, font_color=TEXT_BODY, line_spacing=1.3)
        # 代码位置
        add_rect(slide, x, Inches(4.65), Inches(2.9), Inches(0.3), LIGHT)
        add_text(slide, x, Inches(4.65), Inches(2.9), Inches(0.3),
                 code_ref, font_size=8, font_color=PRIMARY, align="center", anchor="middle", italic=True)
    add_page_chrome(slide, page_num, 20, "踩坑与解决")


def make_testing(prs, page_num):
    """Slide 17: 测试验证"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "03  /  关键技术与对比实验")
    add_slide_title(slide, "测试验证", "7 个自动化测试脚本, 覆盖核心逻辑与数据一致性")

    # 左侧: 大数字
    add_rect(slide, Inches(0.5), Inches(1.75), Inches(3.0), Inches(3.2), PRIMARY)
    add_text(slide, Inches(0.5), Inches(1.85), Inches(3.0), Inches(0.4),
             "测试总览", font_size=12, font_color=LIGHT, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(2.4), Inches(3.0), Inches(1.0),
             "7", font_size=80, font_color=ACCENT, bold=True, align="center", anchor="middle", font_en="Arial")
    add_text(slide, Inches(0.5), Inches(3.5), Inches(3.0), Inches(0.3),
             "自动化测试脚本", font_size=12, font_color=BG, align="center")
    add_text(slide, Inches(0.5), Inches(3.85), Inches(3.0), Inches(0.3),
             "72 项检查", font_size=14, font_color=LIGHT, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(4.2), Inches(3.0), Inches(0.3),
             "全部通过", font_size=14, font_color=SUCCESS, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(4.55), Inches(3.0), Inches(0.3),
             "(每个脚本独立跑, 失败立即报错)",
             font_size=9, font_color=LIGHT, align="center", italic=True)

    # 右侧: 7 个测试列表
    add_text(slide, Inches(3.7), Inches(1.75), Inches(6), Inches(0.3),
             "7 个测试脚本", font_size=14, font_color=PRIMARY, bold=True)
    tests = [
        ("test_preprocess.py",   "7 个原子方法 + 管线调度验证",          "18 项"),
        ("test_add_person.py",   "录入函数去重重构验证",                "9 项"),
        ("test_similarity.py",   "相似度 clamp 到 [0,1] 验证",            "14 项"),
        ("test_sync.py",         "npz / json 同步一致性验证",            "9 项"),
        ("test_ensemble.py",     "EnsembleClassifier 端到端验证",        "22 项"),
        ("test_display_sim.py",  "展示用相似度 helper 验证",             "若干"),
        ("fix_report.py",        "报告改写工具 (辅助脚本)",              "—"),
    ]
    for i, (name, desc, count) in enumerate(tests):
        y = Inches(2.1 + i * 0.38)
        bg_c = LIGHT if i % 2 == 0 else BG
        add_rect(slide, Inches(3.7), y, Inches(5.8), Inches(0.32), bg_c)
        add_text(slide, Inches(3.85), y, Inches(2.0), Inches(0.32),
                 name, font_size=10, font_color=PRIMARY, bold=True, anchor="middle", font_cn="Consolas")
        add_text(slide, Inches(5.9), y, Inches(2.7), Inches(0.32),
                 desc, font_size=10, font_color=TEXT_BODY, anchor="middle")
        add_text(slide, Inches(8.6), y, Inches(0.8), Inches(0.32),
                 count, font_size=10, font_color=SUCCESS, bold=True, align="center", anchor="middle")
    add_page_chrome(slide, page_num, 20, "测试验证")


# ==================== 第 4 章: 总结与展望 ====================
def make_summary(prs, page_num):
    """Slide 19: 项目总结与未来工作"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_label(slide, Inches(0.5), Inches(0.5), "04  /  总结与展望")
    add_slide_title(slide, "项目总结与未来工作", "完成度 + 不足 + 后续可探索的方向")

    # 左侧: 项目成果
    add_text(slide, Inches(0.5), Inches(1.7), Inches(4.5), Inches(0.3),
             "项目成果", font_size=14, font_color=SUCCESS, bold=True)
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(4.5), Inches(2.9), BG,
             line_color=SUCCESS, line_w=Pt(1.5))
    achievements = [
        ("1", "Flask + dlib 端到端 Web 化人脸识别系统, 支持在线录入与实时识别"),
        ("2", "7 步预处理管线统一调度接口, 灵活适配识别 / 录入两种场景"),
        ("3", "13 个 RESTful API 端点, 覆盖人员管理 / 1:1 / 1:N / 特征提取 / 对比实验"),
        ("4", "ResNet 128 维 深度学习 与 4 种手工特征 2329 维 融合的对比实验"),
    ]
    for i, (num, text) in enumerate(achievements):
        y = Inches(2.2 + i * 0.65)
        # 数字圆
        add_rect(slide, Inches(0.7), y, Inches(0.4), Inches(0.4), SUCCESS)
        add_text(slide, Inches(0.7), y, Inches(0.4), Inches(0.4),
                 num, font_size=14, font_color=BG, bold=True, align="center", anchor="middle")
        # 描述
        add_text(slide, Inches(1.2), y, Inches(3.7), Inches(0.6),
                 text, font_size=10, font_color=TEXT_BODY, line_spacing=1.3)

    # 右侧: 不足 + 未来工作
    add_text(slide, Inches(5.2), Inches(1.7), Inches(4.5), Inches(0.3),
             "不足与未来工作", font_size=14, font_color=WARN, bold=True)
    add_rect(slide, Inches(5.2), Inches(2.05), Inches(4.3), Inches(2.9), BG,
             line_color=WARN, line_w=Pt(1.5))
    futures = [
        ("不足 1", "单阈值判定, 没有加权投票 / 多模型融合"),
        ("不足 2", "缺少活体检测, 照片欺骗未防御"),
        ("未来 A", "移动端适配: 微信小程序 / App 拍照录入"),
        ("未来 B", "实时视频流: WebSocket 摄像头实时识别"),
        ("未来 C", "大规模测试: LFW 公开数据集准确率评估"),
        ("未来 D", "模型轻量化: ONNX / TensorRT 部署, 降低推理延迟"),
    ]
    for i, (head, body) in enumerate(futures):
        y = Inches(2.2 + i * 0.43)
        color = DANGER if head.startswith("不足") else ACCENT
        add_text(slide, Inches(5.4), y, Inches(0.7), Inches(0.3),
                 head, font_size=10, font_color=color, bold=True)
        add_text(slide, Inches(6.1), y, Inches(3.3), Inches(0.4),
                 body, font_size=9, font_color=TEXT_BODY)
    add_page_chrome(slide, page_num, 20, "总结与展望")


def make_thanks(prs, page_num):
    """Slide 20: 致谢"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    # 左侧大色块
    add_rect(slide, Inches(0), Inches(0), Inches(3.8), Inches(5.625), PRIMARY)
    add_text(slide, Inches(0.5), Inches(2.0), Inches(3), Inches(0.5),
             "THANKS", font_size=24, font_color=ACCENT, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(2.5), Inches(3), Inches(0.5),
             "致 谢", font_size=40, font_color=BG, bold=True, align="center")
    add_text(slide, Inches(0.5), Inches(3.2), Inches(3), Inches(0.3),
             "感谢聆听", font_size=14, font_color=LIGHT, align="center")
    add_text(slide, Inches(0.5), Inches(3.55), Inches(3), Inches(0.3),
             "请各位老师批评指正", font_size=12, font_color=LIGHT, align="center")
    # 右侧: 关键问题
    add_text(slide, Inches(4.2), Inches(1.4), Inches(5.5), Inches(0.4),
             "老师可能关心的几个问题", font_size=18, font_color=PRIMARY, bold=True)
    add_line(slide, Inches(4.2), Inches(1.9), Inches(5.5), Inches(1.9), ACCENT, 2.0)
    questions = [
        ("Q1", "算法的整体流程是什么?", "8 步: 检测 → 关键点 → 对齐 → 预处理 → ResNet → 128 维 → L2 → 比对"),
        ("Q2", "任意一行代码什么意思?", "L2 归一化 / 距离公式 / 阈值判定 / 录入取平均 — 见代码注释"),
        ("Q3", "性能如何, 能实时吗?", "单张图 ~250ms, 库小识别 ~1ms; 摄像头节流到 100ms / 帧"),
        ("Q4", "创新点在哪?", "统一预处理管线 / 数据一致性 / 深度 vs 手工对比"),
        ("Q5", "实际应用场景?", "门禁 / 课堂签到 / 考勤 — 与任务书选题一致"),
    ]
    for i, (q, head, ans) in enumerate(questions):
        y = Inches(2.05 + i * 0.55)
        add_text(slide, Inches(4.2), y, Inches(0.4), Inches(0.3),
                 q, font_size=11, font_color=ACCENT, bold=True)
        add_text(slide, Inches(4.7), y, Inches(5.0), Inches(0.3),
                 head, font_size=11, font_color=PRIMARY, bold=True)
        add_text(slide, Inches(4.7), y + Inches(0.25), Inches(5.0), Inches(0.3),
                 ans, font_size=9, font_color=TEXT_MUTED)
    add_page_chrome(slide, page_num, 20, "致谢")


# ==================== 主流程 ====================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. 封面
    make_cover(prs)
    # 2. 目录
    make_toc(prs)
    # 3. 第 1 章封面
    make_section_divider(prs, "01", "项目概述", "研究背景 / 核心功能 / 开发历程 / 团队分工 — 从 0 到 1 的故事", 3)
    # 4. 研究背景与意义
    make_bg_meaning(prs, 4)
    # 5. 核心功能与亮点
    make_features(prs, 5)
    # 6. 开发历程
    make_timeline(prs, 6)
    # 7. 团队分工
    make_team(prs, 7)
    # 8. 第 2 章封面
    make_section_divider(prs, "02", "系统架构与核心算法", "5 层模块化架构 + 8 步算法流程 + 7 步预处理管线", 8)
    # 9. 系统整体架构
    make_architecture(prs, 9)
    # 10. 人脸识别算法流程
    make_algorithm_flow(prs, 10)
    # 11. 图像预处理管线
    make_preprocess(prs, 11)
    # 12. 特征提取与匹配
    make_feature_match(prs, 12)
    # 13. 第 3 章封面
    make_section_divider(prs, "03", "关键技术与对比实验", "4 种手工特征 / 深度 vs 手工 / 踩坑与解决 / 测试验证", 13)
    # 14. 4 种手工特征
    make_handcraft(prs, 14)
    # 15. 深度 vs 手工特征
    make_comparison(prs, 15)
    # 16. 踩坑与解决
    make_pitfalls(prs, 16)
    # 17. 测试验证
    make_testing(prs, 17)
    # 18. 第 4 章封面
    make_section_divider(prs, "04", "总结与展望", "项目成果 / 不足反思 / 未来工作方向", 18)
    # 19. 总结与未来工作
    make_summary(prs, 19)
    # 20. 致谢
    make_thanks(prs, 20)

    prs.save(OUT_PATH)
    size_kb = os.path.getsize(OUT_PATH) / 1024
    print(f"Saved: {OUT_PATH}")
    print(f"  Size: {size_kb:.1f} KB")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
